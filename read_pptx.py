from pptx import Presentation

prs = Presentation(r'C:/Users/Administrator/Desktop/产品线/网络协议介绍.pptx')
print(f'总页数: {len(prs.slides)}')
print('\n幻灯片标题:')
for i, slide in enumerate(prs.slides[:15]):
    title = slide.shapes.title.text if slide.shapes.title else "无标题"
    print(f'{i+1}. {title}')

print('\n前8页内容:')
for i, slide in enumerate(prs.slides[:8]):
    print(f'\n=== 第{i+1}页 ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text[:300])
