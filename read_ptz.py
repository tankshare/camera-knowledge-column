from pptx import Presentation

prs = Presentation(r'C:/Users/Administrator/Desktop/产品线/明日小型PTZ产品规划V2.0-2025.07.16.pptx')
print(f'总页数: {len(prs.slides)}')
print('\n幻灯片标题:')
for i, slide in enumerate(prs.slides[:25]):
    title = slide.shapes.title.text if slide.shapes.title else "无标题"
    print(f'{i+1}. {title}')

print('\n前10页详细内容:')
for i, slide in enumerate(prs.slides[:10]):
    print(f'\n=== 第{i+1}页 ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text[:400])
